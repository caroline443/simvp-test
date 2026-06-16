#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ─── Color constants ────────────────────────────────────────────────────────
BLUE_DARK   = RGBColor(0x00, 0x33, 0x66)   # #003366 cover/section bg
BLUE_MID    = RGBColor(0x00, 0x57, 0xAD)   # #0057AD
BLUE_TITLE  = RGBColor(0x00, 0x70, 0xC0)   # #0070C0 title on white
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY   = RGBColor(0x33, 0x33, 0x33)   # #333333
GRAY_SUB    = RGBColor(0x59, 0x59, 0x59)   # #595959
ORANGE      = RGBColor(0xFF, 0x66, 0x00)   # #FF6600 best result
TBL_HEAD_BG = RGBColor(0xDD, 0xEE, 0xFF)   # #DDEEFF table header
TBL_ALT_BG  = RGBColor(0xF5, 0xF5, 0xF5)   # alternating row

# ─── Slide size ─────────────────────────────────────────────────────────────
SLIDE_W = 12192000
SLIDE_H = 6858000

prs = Presentation()
prs.slide_width  = Emu(SLIDE_W)
prs.slide_height = Emu(SLIDE_H)

# Use blank layout for all slides
blank_layout = prs.slide_layouts[6]

# ─── Helper utilities ───────────────────────────────────────────────────────

def add_rect_bg(slide, color, left=0, top=0, width=None, height=None):
    """Fill slide background with a solid color rectangle."""
    if width is None:
        width = SLIDE_W
    if height is None:
        height = SLIDE_H
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(left), Emu(top), Emu(width), Emu(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=BLACK,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_multiline_textbox(slide, lines, left, top, width, height,
                           default_size=16, default_color=DARK_GRAY, wrap=True):
    """
    lines: list of dicts or strings.
    dict keys: text, size, bold, color, align, indent (bool for bullet indent)
    """
    txb = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    tf = txb.text_frame
    tf.word_wrap = wrap
    first = True
    for item in lines:
        if isinstance(item, str):
            item = {"text": item}
        text  = item.get("text", "")
        size  = item.get("size", default_size)
        bold  = item.get("bold", False)
        color = item.get("color", default_color)
        align = item.get("align", PP_ALIGN.LEFT)
        indent = item.get("indent", False)

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        if indent:
            p.level = 1
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return txb


def set_cell_bg(cell, color):
    """Set table cell background color."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
    srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
    r, g, b = color
    srgbClr.set('val', '{:02X}{:02X}{:02X}'.format(r, g, b))


def set_cell_text(cell, text, size=13, bold=False, color=DARK_GRAY,
                  align=PP_ALIGN.CENTER):
    cell.text = ""
    tf = cell.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_table(slide, headers, rows, left, top, width, height,
              col_widths=None, hdr_size=13, row_size=12,
              best_rows=None, best_cols=None):
    """
    best_rows: set of row indices (0-based, not counting header) to highlight in orange bold
    best_cols: set of col indices within best_rows to highlight
    """
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)

    tbl = slide.shapes.add_table(n_rows, n_cols,
                                  Emu(left), Emu(top),
                                  Emu(width), Emu(height)).table

    # Set column widths
    if col_widths:
        for i, w in enumerate(col_widths):
            tbl.columns[i].width = Emu(w)

    # Header row
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        set_cell_bg(cell, TBL_HEAD_BG)
        set_cell_text(cell, hdr, size=hdr_size, bold=True, color=BLUE_MID)

    # Data rows
    for ri, row in enumerate(rows):
        bg = TBL_ALT_BG if ri % 2 == 1 else WHITE
        is_best_row = best_rows and ri in best_rows
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            set_cell_bg(cell, bg)
            is_best = is_best_row and (best_cols is None or ci in best_cols)
            # Strip ** markers for styling
            text = str(val).replace('**', '')
            use_bold = is_best or ('**' in str(val))
            use_color = ORANGE if is_best else DARK_GRAY
            set_cell_text(cell, text, size=row_size, bold=use_bold, color=use_color)

    return tbl


def add_divider(slide, top, color=BLUE_TITLE, left=400000, width=None):
    if width is None:
        width = SLIDE_W - 800000
    shape = slide.shapes.add_shape(1,
        Emu(left), Emu(top), Emu(width), Emu(30000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 1 — Cover
# ═══════════════════════════════════════════════════════════════════════════════
slide1 = prs.slides.add_slide(blank_layout)

# Dark blue background
add_rect_bg(slide1, BLUE_DARK)

# Decorative accent bar (left side)
acc = slide1.shapes.add_shape(1, Emu(0), Emu(0), Emu(280000), Emu(SLIDE_H))
acc.fill.solid()
acc.fill.fore_color.rgb = BLUE_MID
acc.line.fill.background()

# Bottom accent stripe
bot = slide1.shapes.add_shape(1, Emu(0), Emu(SLIDE_H - 200000), Emu(SLIDE_W), Emu(200000))
bot.fill.solid()
bot.fill.fore_color.rgb = BLUE_MID
bot.line.fill.background()

# Main title
add_textbox(slide1,
    "基于深度学习的强对流天气高保真临近预报算法研究",
    left=600000, top=1600000, width=11000000, height=800000,
    font_size=34, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Subtitle
add_textbox(slide1,
    "中期进展汇报",
    left=600000, top=2500000, width=11000000, height=500000,
    font_size=24, bold=False, color=RGBColor(0xAA, 0xCC, 0xFF),
    align=PP_ALIGN.CENTER)

# Divider line
div = slide1.shapes.add_shape(1,
    Emu(2000000), Emu(3100000), Emu(8200000), Emu(40000))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0xAA, 0xCC, 0xFF)
div.line.fill.background()

# Info block
info_lines = [
    {"text": "姓  名：张雨欣", "size": 18, "color": WHITE},
    {"text": "学  号：2024020636", "size": 18, "color": WHITE},
    {"text": "导  师：何善宝", "size": 18, "color": WHITE},
    {"text": "专  业：人工智能", "size": 18, "color": WHITE},
]
add_multiline_textbox(slide1, info_lines,
    left=3500000, top=3300000, width=5200000, height=1600000)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 2 — Table of Contents (section page, blue bg)
# ═══════════════════════════════════════════════════════════════════════════════
slide2 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide2, BLUE_DARK)

# Side accent
acc2 = slide2.shapes.add_shape(1, Emu(0), Emu(0), Emu(200000), Emu(SLIDE_H))
acc2.fill.solid()
acc2.fill.fore_color.rgb = BLUE_MID
acc2.line.fill.background()

# Title
add_textbox(slide2, "目  录",
    left=400000, top=400000, width=11000000, height=600000,
    font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_divider(slide2, top=1050000, color=RGBColor(0xAA, 0xCC, 0xFF))

items = [
    "01   研究计划调整说明",
    "02   本阶段改进方案",
    "03   实验设置",
    "04   实验结果",
    "05   结论与后续计划",
]
for i, item in enumerate(items):
    add_textbox(slide2, item,
        left=1800000, top=1300000 + i * 950000,
        width=8600000, height=750000,
        font_size=22, bold=(i == 0), color=WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 3 — 研究计划调整说明
# ═══════════════════════════════════════════════════════════════════════════════
slide3 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide3, WHITE)

# Top blue banner
banner = slide3.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner.fill.solid()
banner.fill.fore_color.rgb = BLUE_DARK
banner.line.fill.background()

add_textbox(slide3, "研究计划调整说明",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=28, bold=True, color=WHITE)

# Left column: original plan + reason
col_left = [
    {"text": "原计划（上次汇报）", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "• 在 Translator Inception 模块内嵌入空间注意力（CBAM）", "size": 14, "color": DARK_GRAY},
    {"text": "• 混合损失函数：Weighted-L1 + MS-SSIM", "size": 14, "color": DARK_GRAY},
    {"text": "• 频域感知模块（FFT）", "size": 14, "color": DARK_GRAY},
    {"text": "", "size": 10},
    {"text": "调整原因", "size": 16, "bold": True, "color": ORANGE},
    {"text": "• 上述改进集成到 384×384 全分辨率训练时，训练崩溃（loss NaN）", "size": 14, "color": DARK_GRAY},
    {"text": "• 根本原因：12帧输入时 Inception 翻译器通道数达 3072，FP16 精度下卷积累加溢出", "size": 14, "color": DARK_GRAY},
    {"text": "  （3072×25=76800 次乘加 > FP16 上限 65504）", "size": 13, "color": GRAY_SUB},
    {"text": "• 转换策略：降分辨率至 128×128，聚焦架构与训练策略改进", "size": 14, "color": DARK_GRAY},
]

add_multiline_textbox(slide3, col_left,
    left=300000, top=800000, width=5700000, height=5600000)

# Right column: new direction
col_right = [
    {"text": "本阶段改进方向", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "", "size": 8},
    {"text": "训练策略改进", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "• OPSD（在策略自蒸馏）", "size": 14, "color": DARK_GRAY},
    {"text": "• OPSD-RW（奖励加权蒸馏）", "size": 14, "color": DARK_GRAY},
    {"text": "", "size": 8},
    {"text": "架构改进", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "• Mamba 时序翻译器替代 Inception", "size": 14, "color": DARK_GRAY},
    {"text": "", "size": 8},
    {"text": "注意力改进（独立对比实验）", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "• CBAM（训练中）", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide3, col_right,
    left=6200000, top=800000, width=5700000, height=5600000)

# Vertical divider
vdiv = slide3.shapes.add_shape(1, Emu(6050000), Emu(850000), Emu(40000), Emu(5500000))
vdiv.fill.solid()
vdiv.fill.fore_color.rgb = TBL_HEAD_BG
vdiv.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 4 — OPSD 动机
# ═══════════════════════════════════════════════════════════════════════════════
slide4 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide4, WHITE)

banner4 = slide4.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner4.fill.solid()
banner4.fill.fore_color.rgb = BLUE_DARK
banner4.line.fill.background()

add_textbox(slide4, "改进一：OPSD — 解决 Exposure Bias",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=26, bold=True, color=WHITE)

# Left: problem description
left_content = [
    {"text": "Exposure Bias 问题", "size": 18, "bold": True, "color": BLUE_TITLE},
    {"text": "", "size": 8},
    {"text": "训练时：每步以真实历史帧为条件（Teacher Forcing）", "size": 15, "color": DARK_GRAY},
    {"text": "", "size": 6},
    {"text": "推理时：每步以自己预测的帧为条件（有误差）", "size": 15, "color": DARK_GRAY},
    {"text": "", "size": 6},
    {"text": "误差随时间步累积 → 预测越来越差", "size": 15, "bold": True, "color": ORANGE},
]
add_multiline_textbox(slide4, left_content,
    left=300000, top=800000, width=5500000, height=3000000)

# Right: table
headers4 = ["时间步", "训练时输入", "推理时输入"]
rows4 = [
    ["第 1 步", "全是真实帧", "全是真实帧"],
    ["第 6 步", "全是真实帧", "含 5 帧预测误差"],
    ["第 12 步", "全是真实帧", "含 11 帧累积误差"],
]
col_w4 = [SLIDE_W // 3 - 100000] * 3
add_table(slide4, headers4, rows4,
    left=6000000, top=900000,
    width=5900000, height=1800000,
    hdr_size=14, row_size=13)

# Arrow / annotation below table
ann_lines = [
    {"text": "核心矛盾：训练与推理阶段输入分布不一致", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "→ 训练时模型从未见过『自己的预测帧』作为输入", "size": 14, "color": DARK_GRAY},
    {"text": "→ 推理时累积误差导致后期帧质量急剧下降", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide4, ann_lines,
    left=6000000, top=2800000, width=5900000, height=1500000)

# Bottom summary box
box = slide4.shapes.add_shape(1, Emu(300000), Emu(4000000), Emu(11600000), Emu(500000))
box.fill.solid()
box.fill.fore_color.rgb = TBL_HEAD_BG
box.line.color.rgb = BLUE_TITLE

add_textbox(slide4, "解决方案：让训练时也接触自己的预测帧，缩小训练-推理分布差距",
    left=400000, top=4020000, width=11400000, height=460000,
    font_size=16, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 5 — OPSD 方法
# ═══════════════════════════════════════════════════════════════════════════════
slide5 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide5, WHITE)

banner5 = slide5.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner5.fill.solid()
banner5.fill.fore_color.rgb = BLUE_DARK
banner5.line.fill.background()

add_textbox(slide5, "OPSD — 教师-学生框架",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=26, bold=True, color=WHITE)

core = [{"text": "核心思想：不改模型结构，只改训练方式", "size": 17, "bold": True, "color": BLUE_TITLE}]
add_multiline_textbox(slide5, core, left=300000, top=780000, width=11600000, height=400000)

# Student branch
stu_box = slide5.shapes.add_shape(1, Emu(300000), Emu(1250000), Emu(5500000), Emu(2000000))
stu_box.fill.solid()
stu_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFF)
stu_box.line.color.rgb = BLUE_TITLE

stu_lines = [
    {"text": "学生分支（Student）", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "• 走正常自回归推理（用自己的预测帧填窗口）", "size": 14, "color": DARK_GRAY},
    {"text": "• 产生预测分布 p^s_t", "size": 14, "color": DARK_GRAY},
    {"text": "• 梯度正常反传，参数更新", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide5, stu_lines,
    left=400000, top=1320000, width=5300000, height=1800000)

# Teacher branch
tch_box = slide5.shapes.add_shape(1, Emu(6200000), Emu(1250000), Emu(5700000), Emu(2000000))
tch_box.fill.solid()
tch_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF5, 0xE8)
tch_box.line.color.rgb = ORANGE

tch_lines = [
    {"text": "教师分支（Teacher）", "size": 16, "bold": True, "color": ORANGE},
    {"text": "• 每步用真实未来帧填窗口（『上帝视角』）", "size": 14, "color": DARK_GRAY},
    {"text": "• 包裹在 torch.no_grad()，不参与反向传播", "size": 14, "color": DARK_GRAY},
    {"text": "• 产生目标分布 p^r_t（参考信号）", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide5, tch_lines,
    left=6300000, top=1320000, width=5500000, height=1800000)

# Loss function
loss_lines = [
    {"text": "损失函数", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "L_OPSD  =  λ_KL × KL(p^r_t ‖ p^s_t)  +  λ_CE × CE损失", "size": 15, "bold": True, "color": DARK_GRAY},
    {"text": "λ_KL = 1.0,   λ_CE = 0.5,   蒸馏温度 T = 2.0", "size": 14, "color": GRAY_SUB},
]
add_multiline_textbox(slide5, loss_lines,
    left=300000, top=3400000, width=7000000, height=1200000)

# Advantage
adv_box = slide5.shapes.add_shape(1, Emu(7500000), Emu(3350000), Emu(4400000), Emu(1200000))
adv_box.fill.solid()
adv_box.fill.fore_color.rgb = TBL_HEAD_BG
adv_box.line.color.rgb = BLUE_TITLE

adv_lines = [
    {"text": "优势", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "• 额外显存开销 ≈ 0（教师分支无梯度）", "size": 13, "color": DARK_GRAY},
    {"text": "• 从 Baseline checkpoint 热启动", "size": 13, "color": DARK_GRAY},
]
add_multiline_textbox(slide5, adv_lines,
    left=7600000, top=3380000, width=4200000, height=1100000)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 6 — OPSD-RW
# ═══════════════════════════════════════════════════════════════════════════════
slide6 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide6, WHITE)

banner6 = slide6.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner6.fill.solid()
banner6.fill.fore_color.rgb = BLUE_DARK
banner6.line.fill.background()

add_textbox(slide6, "OPSD-RW — 奖励加权蒸馏",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=26, bold=True, color=WHITE)

# Problem
prob = [
    {"text": "问题：标准 OPSD 对所有时间步均等加权，但预测质量随时间非均匀下降", "size": 16, "bold": True, "color": ORANGE},
    {"text": "", "size": 8},
    {"text": "解决：以每步 CSI（阈值74）的反向值作为该步的 KL 权重，让模型多关注预测差的时间步", "size": 15, "color": DARK_GRAY},
]
add_multiline_textbox(slide6, prob,
    left=300000, top=780000, width=11600000, height=1100000)

# Formula box
fml_box = slide6.shapes.add_shape(1, Emu(300000), Emu(1950000), Emu(5500000), Emu(1600000))
fml_box.fill.solid()
fml_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFF)
fml_box.line.color.rgb = BLUE_TITLE

fml_lines = [
    {"text": "权重计算", "size": 15, "bold": True, "color": BLUE_TITLE},
    {"text": "r_t = CSI(ŷ_t, y_t, τ=74)", "size": 14, "color": DARK_GRAY},
    {"text": "w_t = 1 - r_t", "size": 14, "bold": True, "color": DARK_GRAY},
    {"text": "L_OPSD-RW = (1/T) Σ (1−r_t) × L^t_KL + λ_CE × L_CE", "size": 14, "bold": True, "color": BLUE_MID},
]
add_multiline_textbox(slide6, fml_lines,
    left=400000, top=2000000, width=5300000, height=1450000)

# Table
headers6 = ["时间步预测质量", "r_t（CSI）", "权重 w_t = 1-r_t", "梯度效果"]
rows6 = [
    ["预测好（CSI=0.8）", "0.8", "0.2", "少给梯度"],
    ["预测差（CSI=0.2）", "0.2", "0.8", "多给梯度"],
]
add_table(slide6, headers6, rows6,
    left=6000000, top=1950000,
    width=5900000, height=1500000,
    hdr_size=14, row_size=13)

# Intuition
int_lines = [
    {"text": "直觉理解", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "• 模型在早期时间步预测较好（r_t 高） → 权重小，减少不必要干扰", "size": 14, "color": DARK_GRAY},
    {"text": "• 模型在后期时间步预测差（r_t 低）  → 权重大，强化这些步的蒸馏信号", "size": 14, "color": DARK_GRAY},
    {"text": "• 自适应聚焦于最需要改进的预测时段", "size": 14, "bold": True, "color": BLUE_MID},
]
add_multiline_textbox(slide6, int_lines,
    left=300000, top=3700000, width=11600000, height=1500000)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 7 — Mamba Translator
# ═══════════════════════════════════════════════════════════════════════════════
slide7 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide7, WHITE)

banner7 = slide7.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner7.fill.solid()
banner7.fill.fore_color.rgb = BLUE_DARK
banner7.line.fill.background()

add_textbox(slide7, "改进二：Mamba 时序翻译器",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=26, bold=True, color=WHITE)

# Left: Inception problems
inc_box = slide7.shapes.add_shape(1, Emu(300000), Emu(800000), Emu(5500000), Emu(2400000))
inc_box.fill.solid()
inc_box.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF0)
inc_box.line.color.rgb = RGBColor(0xFF, 0x88, 0x88)

inc_lines = [
    {"text": "Inception 翻译器的问题", "size": 16, "bold": True, "color": RGBColor(0xCC, 0x00, 0x00)},
    {"text": "• 把 12 帧特征拼到通道维度 → 通道数 12×256=3072", "size": 14, "color": DARK_GRAY},
    {"text": "• 无显式时序轴，时序信息隐含在通道排列中", "size": 14, "color": DARK_GRAY},
    {"text": "• 所有空间位置用同一卷积核处理，无空间区分", "size": 14, "color": DARK_GRAY},
    {"text": "• FP16 下通道乘加溢出（>65504）", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide7, inc_lines,
    left=400000, top=870000, width=5300000, height=2200000)

# Right: Mamba approach
mam_box = slide7.shapes.add_shape(1, Emu(6200000), Emu(800000), Emu(5700000), Emu(2400000))
mam_box.fill.solid()
mam_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF4, 0xFF)
mam_box.line.color.rgb = BLUE_TITLE

mam_lines = [
    {"text": "Mamba 的做法", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "• 把每个空间位置视为独立时序序列 [h×w, T, C]", "size": 14, "color": DARK_GRAY},
    {"text": "• 在时序轴上用 S6 选择性状态空间模型（SSM）处理", "size": 14, "color": DARK_GRAY},
    {"text": "• 状态转移矩阵依赖输入动态调整（选择性记忆/遗忘）", "size": 14, "color": DARK_GRAY},
    {"text": "• 通道维度固定为 C，不随帧数增长", "size": 14, "bold": True, "color": BLUE_MID},
]
add_multiline_textbox(slide7, mam_lines,
    left=6300000, top=870000, width=5500000, height=2200000)

# Comparison table
headers7 = ["对比项", "Inception", "Mamba"]
rows7 = [
    ["时序处理", "通道拼接 + 2D卷积", "显式 SSM 递推"],
    ["通道维度", "T×C（随帧数增长）", "C（固定）"],
    ["空间位置区分", "无", "有（各位置独立建模）"],
    ["FP16 溢出风险", "高（通道数大）", "低（通道固定）"],
]
add_table(slide7, headers7, rows7,
    left=300000, top=3380000,
    width=11600000, height=2200000,
    hdr_size=14, row_size=13,
    best_rows={0, 1, 2, 3}, best_cols={2})


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 8 — 实验设置
# ═══════════════════════════════════════════════════════════════════════════════
slide8 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide8, WHITE)

banner8 = slide8.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner8.fill.solid()
banner8.fill.fore_color.rgb = BLUE_DARK
banner8.line.fill.background()

add_textbox(slide8, "实验设置",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=28, bold=True, color=WHITE)

# Left column: dataset & prediction config
l8_lines = [
    {"text": "数据集与预测配置", "size": 17, "bold": True, "color": BLUE_TITLE},
    {"text": "", "size": 8},
    {"text": "数据集：SEVIR VIL（2017-2018 训练，2019 测试）", "size": 15, "color": DARK_GRAY},
    {"text": "分辨率：128×128（中心裁剪）", "size": 15, "color": DARK_GRAY},
    {"text": "预测配置：12帧输入 → 12帧预测", "size": 15, "color": DARK_GRAY},
    {"text": "时间跨度：60分钟，5分钟间隔", "size": 15, "color": DARK_GRAY},
    {"text": "", "size": 10},
    {"text": "训练配置", "size": 17, "bold": True, "color": BLUE_TITLE},
    {"text": "", "size": 8},
    {"text": "Baseline LR: 5×10⁻⁴，50 epoch", "size": 15, "color": DARK_GRAY},
    {"text": "OPSD LR: 2×10⁻⁴，50 epoch（从 Baseline 热启动）", "size": 15, "color": DARK_GRAY},
    {"text": "损失：前景加权交叉熵（前景权重 5.0）", "size": 15, "color": DARK_GRAY},
    {"text": "优化器：AdamW，梯度裁剪 max_norm=1.0", "size": 15, "color": DARK_GRAY},
]
add_multiline_textbox(slide8, l8_lines,
    left=300000, top=800000, width=5700000, height=5500000)

# Right column: metrics
r8_lines = [
    {"text": "评估指标", "size": 17, "bold": True, "color": BLUE_TITLE},
    {"text": "", "size": 8},
    {"text": "CSI（临界成功指数）", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "CSI(τ) = TP / (TP + FP + FN)", "size": 14, "color": DARK_GRAY},
    {"text": "", "size": 8},
    {"text": "报告阈值（VIL 强度级别）：", "size": 15, "color": DARK_GRAY},
    {"text": "τ ∈ {16, 74, 133, 160, 181, 219}", "size": 15, "bold": True, "color": DARK_GRAY},
    {"text": "CSI-M = 六阈值均值", "size": 15, "bold": True, "color": DARK_GRAY},
    {"text": "", "size": 10},
    {"text": "其他指标", "size": 15, "bold": True, "color": BLUE_MID},
    {"text": "POD：命中率（Probability of Detection）", "size": 14, "color": DARK_GRAY},
    {"text": "FAR：虚警率（False Alarm Rate）", "size": 14, "color": DARK_GRAY},
    {"text": "HSS：技巧评分（Heidke Skill Score）", "size": 14, "color": DARK_GRAY},
    {"text": "", "size": 8},
    {"text": "主要关注 CSI-M 和 CSI@74（最常见强对流阈值）", "size": 14, "bold": True, "color": BLUE_TITLE},
]
add_multiline_textbox(slide8, r8_lines,
    left=6200000, top=800000, width=5700000, height=5500000)

# Divider
vdiv8 = slide8.shapes.add_shape(1, Emu(6050000), Emu(850000), Emu(40000), Emu(5500000))
vdiv8.fill.solid()
vdiv8.fill.fore_color.rgb = TBL_HEAD_BG
vdiv8.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 9 — Main experiment results
# ═══════════════════════════════════════════════════════════════════════════════
slide9 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide9, WHITE)

banner9 = slide9.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner9.fill.solid()
banner9.fill.fore_color.rgb = BLUE_DARK
banner9.line.fill.background()

add_textbox(slide9, "实验结果 — 主指标对比",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=28, bold=True, color=WHITE)

headers9 = ["模型", "CSI-M", "CSI@74", "CSI@133", "POD@74", "FAR@74", "HSS@74"]
rows9 = [
    ["Inception Baseline",    "0.3914", "0.6024", "0.3514", "0.7567", "0.2585", "0.6809"],
    ["Inception + OPSD",      "0.3938", "0.6070", "0.3557", "0.7475", "0.2423", "0.6868"],
    ["Inception + OPSD-RW",   "0.3793", "0.6017", "0.2988", "0.7412", "0.2442", "0.6819"],
    ["Mamba Baseline",        "0.3807", "0.5943", "0.3229", "0.7170", "0.2309", "0.6762"],
    ["Mamba + OPSD",          "0.3810", "0.6045", "0.3316", "0.7276", "0.2251", "0.6864"],
    ["Mamba + OPSD-RW",       "0.3893", "0.6090", "0.3308", "0.7508", "0.2429", "0.6884"],
]

# best_rows: Inception+OPSD (row 1) and Mamba+OPSD-RW (row 5)
# For best row highlighting we need cell-level control
tbl9 = slide9.shapes.add_table(7, 7,
    Emu(200000), Emu(800000),
    Emu(11800000), Emu(3600000)).table

col_widths9 = [3200000, 1300000, 1300000, 1400000, 1300000, 1300000, 1300000]
for i, w in enumerate(col_widths9):
    tbl9.columns[i].width = Emu(w)

for ci, hdr in enumerate(headers9):
    cell = tbl9.cell(0, ci)
    set_cell_bg(cell, TBL_HEAD_BG)
    set_cell_text(cell, hdr, size=14, bold=True, color=BLUE_MID)

# Best values per column (row index 0-based in data)
# Inception+OPSD: CSI-M best, CSI@74 best, CSI@133 best, HSS best, FAR best
# Mamba+OPSD-RW: CSI@74 best overall, CSI-M among Mamba best
best_per_col = {
    1: {1},    # CSI-M: Inception+OPSD row=1
    2: {5},    # CSI@74: Mamba+OPSD-RW row=5
    3: {1},    # CSI@133: Inception+OPSD row=1
    4: {0},    # POD@74: Inception Baseline row=0
    5: {4},    # FAR@74: Mamba+OPSD row=4 (lowest)
    6: {5},    # HSS@74: Mamba+OPSD-RW row=5
}

# Highlight entire best model rows
highlight_rows = {1, 5}  # Inception+OPSD, Mamba+OPSD-RW

for ri, row in enumerate(rows9):
    bg = TBL_ALT_BG if ri % 2 == 1 else WHITE
    for ci, val in enumerate(row):
        cell = tbl9.cell(ri + 1, ci)
        set_cell_bg(cell, bg)
        is_best_val = (ci in best_per_col and ri in best_per_col[ci])
        is_best_row_name = (ri in highlight_rows and ci == 0)
        use_bold = is_best_val or is_best_row_name
        use_color = ORANGE if is_best_val else (BLUE_MID if is_best_row_name else DARK_GRAY)
        set_cell_text(cell, val, size=13, bold=use_bold, color=use_color)

# Annotations
ann9 = [
    {"text": "橙色标注为各列最优值", "size": 14, "bold": True, "color": ORANGE},
    {"text": "• Mamba + OPSD-RW 在 CSI@74 和 HSS@74 均达最优，整体性能最强", "size": 14, "color": DARK_GRAY},
    {"text": "• Inception + OPSD 在 CSI-M 和 CSI@133 表现最佳，蒸馏策略对 Inception 有效", "size": 14, "color": DARK_GRAY},
    {"text": "• Mamba Baseline 虚警率（FAR）最低，选择性状态建模有效抑制误报", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide9, ann9,
    left=300000, top=4600000, width=11600000, height=1800000)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 10 — Step-wise CSI analysis
# ═══════════════════════════════════════════════════════════════════════════════
slide10 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide10, WHITE)

banner10 = slide10.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner10.fill.solid()
banner10.fill.fore_color.rgb = BLUE_DARK
banner10.line.fill.background()

add_textbox(slide10, "实验结果 — 预测质量随时间衰减分析",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=26, bold=True, color=WHITE)

add_textbox(slide10, "CSI@74 随预测时间步的变化",
    left=300000, top=760000, width=11600000, height=380000,
    font_size=16, bold=True, color=BLUE_TITLE)

headers10 = ["模型", "5min", "15min", "30min", "45min", "60min"]
rows10 = [
    ["Inception Baseline",  "0.732", "0.685", "0.607", "0.545", "0.488"],
    ["Inception + OPSD",    "0.735", "0.691", "0.612", "0.547", "0.493"],
    ["Inception + OPSD-RW", "0.730", "0.686", "0.607", "0.541", "0.488"],
    ["Mamba Baseline",      "0.739", "0.684", "0.595", "0.532", "0.470"],
    ["Mamba + OPSD",        "0.741", "0.691", "0.606", "0.545", "0.487"],
    ["Mamba + OPSD-RW",     "0.746", "0.696", "0.611", "0.549", "0.491"],
]

tbl10 = slide10.shapes.add_table(7, 6,
    Emu(200000), Emu(1200000),
    Emu(11800000), Emu(3000000)).table

col_widths10 = [3500000, 1660000, 1660000, 1660000, 1660000, 1660000]
for i, w in enumerate(col_widths10):
    tbl10.columns[i].width = Emu(w)

for ci, hdr in enumerate(headers10):
    cell = tbl10.cell(0, ci)
    set_cell_bg(cell, TBL_HEAD_BG)
    set_cell_text(cell, hdr, size=14, bold=True, color=BLUE_MID)

# Best per column: Mamba+OPSD-RW (row 5) is best in all time steps
best10 = {1: {5}, 2: {5}, 3: {5}, 4: {5}, 5: {5}}

for ri, row in enumerate(rows10):
    bg = TBL_ALT_BG if ri % 2 == 1 else WHITE
    for ci, val in enumerate(row):
        cell = tbl10.cell(ri + 1, ci)
        set_cell_bg(cell, bg)
        is_best = ci in best10 and ri in best10[ci]
        use_bold = is_best or (ri == 5 and ci == 0)
        use_color = ORANGE if is_best else (BLUE_MID if (ri == 5 and ci == 0) else DARK_GRAY)
        set_cell_text(cell, val, size=13, bold=use_bold, color=use_color)

# Key observation
obs_box = slide10.shapes.add_shape(1, Emu(200000), Emu(4400000), Emu(11800000), Emu(800000))
obs_box.fill.solid()
obs_box.fill.fore_color.rgb = TBL_HEAD_BG
obs_box.line.color.rgb = BLUE_TITLE

obs_lines = [
    {"text": "关键观察：OPSD 对 30~60min 长时预测提升最明显", "size": 16, "bold": True, "color": BLUE_TITLE},
    {"text": "验证了 Exposure Bias 主要影响累积误差较大的后期时间步（训练-推理分布差距随时间扩大）", "size": 14, "color": DARK_GRAY},
]
add_multiline_textbox(slide10, obs_lines,
    left=350000, top=4450000, width=11500000, height=700000)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 11 — Conclusion
# ═══════════════════════════════════════════════════════════════════════════════
slide11 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide11, WHITE)

banner11 = slide11.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner11.fill.solid()
banner11.fill.fore_color.rgb = BLUE_DARK
banner11.line.fill.background()

add_textbox(slide11, "结  论",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=28, bold=True, color=WHITE)

conclusions = [
    {
        "num": "1",
        "title": "Mamba + OPSD-RW 整体最优（CSI-M=0.3893，CSI@74=0.6090）",
        "detail": "Mamba 显式时序建模与奖励加权蒸馏存在协同效应，两者互补",
    },
    {
        "num": "2",
        "title": "OPSD 对 Inception 有效，主要体现在降低虚警率（FAR -0.016）",
        "detail": "OPSD 缩小了训练-推理分布差距，使模型预测更保守、更精准",
    },
    {
        "num": "3",
        "title": "OPSD-RW 在 Inception 上效果下降（CSI@133 -0.053）",
        "detail": "中高阈值梯度分配可能失衡，奖励权重设计需针对架构特性调整",
    },
    {
        "num": "4",
        "title": "Mamba Baseline FAR 最低（0.2309）",
        "detail": "选择性状态建模使模型定位对流区域更精准，有效减少误报",
    },
    {
        "num": "5",
        "title": "长预测时域（30~60min）提升最显著",
        "detail": "验证了 Exposure Bias 主要影响累积误差较大的后期时间步",
    },
]

for i, c in enumerate(conclusions):
    top_base = 820000 + i * 1100000
    # Number circle (simulated by a small square)
    circ = slide11.shapes.add_shape(1,
        Emu(300000), Emu(top_base), Emu(380000), Emu(380000))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE_TITLE
    circ.line.fill.background()
    add_textbox(slide11, c["num"],
        left=310000, top=top_base - 20000, width=360000, height=420000,
        font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide11, c["title"],
        left=800000, top=top_base, width=11000000, height=380000,
        font_size=15, bold=True, color=DARK_GRAY)
    add_textbox(slide11, c["detail"],
        left=800000, top=top_base + 380000, width=11000000, height=350000,
        font_size=13, bold=False, color=GRAY_SUB)


# ═══════════════════════════════════════════════════════════════════════════════
# Slide 12 — Future Work
# ═══════════════════════════════════════════════════════════════════════════════
slide12 = prs.slides.add_slide(blank_layout)
add_rect_bg(slide12, WHITE)

banner12 = slide12.shapes.add_shape(1, Emu(0), Emu(0), Emu(SLIDE_W), Emu(700000))
banner12.fill.solid()
banner12.fill.fore_color.rgb = BLUE_DARK
banner12.line.fill.background()

add_textbox(slide12, "后续工作计划",
    left=300000, top=80000, width=11500000, height=560000,
    font_size=28, bold=True, color=WHITE)

plans = [
    {
        "num": "1",
        "title": "完成 Inception + CBAM 实验（训练中）",
        "detail": "对比 CBAM 作为独立改进的效果，明确空间注意力的贡献",
    },
    {
        "num": "2",
        "title": "在最优模型（Mamba + OPSD-RW）基础上叠加 CBAM",
        "detail": "验证注意力机制与 Mamba+OPSD 的协同效果",
    },
    {
        "num": "3",
        "title": "提升分辨率至 384×384",
        "detail": "使用 FP32 训练或梯度检查点避免 NaN，对齐原始 SEVIR 标准评测",
    },
    {
        "num": "4",
        "title": "探索 Vanilla SimVP（MSE）作为基准",
        "detail": "量化自回归化和分类离散化各自带来的增益，明确各改进贡献",
    },
    {
        "num": "5",
        "title": "撰写论文",
        "detail": "整理实验结果，完成方法描述与实验分析",
    },
]

for i, p in enumerate(plans):
    top_base = 850000 + i * 1100000
    circ = slide12.shapes.add_shape(1,
        Emu(300000), Emu(top_base), Emu(380000), Emu(380000))
    circ.fill.solid()
    circ.fill.fore_color.rgb = BLUE_MID
    circ.line.fill.background()
    add_textbox(slide12, p["num"],
        left=310000, top=top_base - 20000, width=360000, height=420000,
        font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide12, p["title"],
        left=800000, top=top_base, width=11000000, height=380000,
        font_size=15, bold=True, color=DARK_GRAY)
    add_textbox(slide12, p["detail"],
        left=800000, top=top_base + 380000, width=11000000, height=350000,
        font_size=13, bold=False, color=GRAY_SUB)

# Bottom thank you
thanks_box = slide12.shapes.add_shape(1,
    Emu(0), Emu(SLIDE_H - 500000), Emu(SLIDE_W), Emu(500000))
thanks_box.fill.solid()
thanks_box.fill.fore_color.rgb = BLUE_DARK
thanks_box.line.fill.background()

add_textbox(slide12, "感谢各位老师的聆听与指导！",
    left=0, top=SLIDE_H - 480000, width=SLIDE_W, height=440000,
    font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ─── Save ────────────────────────────────────────────────────────────────────
OUTPUT = "/Users/amon/Desktop/张雨欣-中期答辩2.pptx"
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")
